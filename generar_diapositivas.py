# -*- coding: utf-8 -*-
"""
============================================================================
 GENERADOR DE DIAPOSITIVAS — INFORME MENSUAL DE PAGOS Y CTS
============================================================================
Arma automaticamente 3 diapositivas de reporte (KPIs 1/2, KPIs 2/2 y Gestion
Operativa) a partir de los KPIs calculados por "generar_kpis.py": graficos
nativos de PowerPoint, tarjetas de "que significa" en lenguaje simple y
tarjetas de indicadores clave.

Esta version es de codigo abierto / portafolio: construye la presentacion
desde una diapositiva en blanco (no depende de ninguna plantilla corporativa)
y usa una paleta de marca de EJEMPLO. Para usarla en tu empresa, cambia la
paleta (PALETA DE MARCA) y, si quieres, agrega tu propio logo con LOGO_PNG.

Uso: python generar_diapositivas.py
Salida: "Informe_Demo_KPIs.pptx"
============================================================================
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import generar_kpis as g

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn

CARPETA = os.path.dirname(os.path.abspath(__file__))
SALIDA = os.path.join(CARPETA, "Informe_Demo_KPIs.pptx")

# Logo opcional: si el archivo existe se coloca en la esquina del banner; si
# no existe, el banner se genera igual pero sin logo.
LOGO_PNG = os.path.join(CARPETA, "logo_ejemplo.png")

# ---------------------------------------------------------------------------
# PALETA DE MARCA (EJEMPLO — reemplaza estos colores por los de tu empresa)
# ---------------------------------------------------------------------------
ROJO       = RGBColor(0xE2, 0x00, 0x0C)
NARANJA    = RGBColor(0xE8, 0x42, 0x04)
GRIS_OSC   = RGBColor(0x40, 0x40, 0x40)
GRIS_MED   = RGBColor(0x70, 0x70, 0x70)
VERDE      = RGBColor(0x4C, 0xAF, 0x50)
BLANCO     = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO      = RGBColor(0x20, 0x20, 0x20)
FONDO_GRIS = RGBColor(0xF5, 0xF5, 0xF5)

FUENTE = "Calibri"
MESES_LABEL = {"2026-04": "Abr-26", "2026-05": "May-26", "2026-06": "Jun-26", "2026-07": "Jul-26"}

# Tamaño de diapositiva 16:9 (ancho, alto)
SLIDE_W, SLIDE_H = Cm(33.87), Cm(19.05)


# ---------------------------------------------------------------------------
# HELPERS DE FORMATO
# ---------------------------------------------------------------------------
def _set_no_line(shape):
    shape.line.fill.background()


def _texto(tf, texto, size=11, bold=False, color=NEGRO, align=PP_ALIGN.LEFT, font=FUENTE):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = texto
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return p


def _add_textbox(slide, x, y, w, h, texto, size=11, bold=False, color=NEGRO,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    _texto(tf, texto, size=size, bold=bold, color=color, align=align)
    return box


def banner_titulo(slide, titulo, subtitulo=None):
    """Banner rojo curvo superior izquierdo + logo opcional a la derecha."""
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(0), Cm(0), Cm(20), Cm(1.7))
    banner.adjustments[0] = 0.5
    banner.fill.solid(); banner.fill.fore_color.rgb = ROJO
    _set_no_line(banner); banner.shadow.inherit = False
    tf = banner.text_frame
    tf.margin_left = Cm(0.6); tf.margin_top = Cm(0.15); tf.margin_bottom = Cm(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _texto(tf, titulo, size=20, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)
    if os.path.exists(LOGO_PNG):
        slide.shapes.add_picture(LOGO_PNG, Cm(27.3), Cm(0.35), height=Cm(1.0))
    if subtitulo:
        _add_textbox(slide, Cm(0.3), Cm(1.85), Cm(24), Cm(0.6), subtitulo, size=11, color=GRIS_MED)


def tarjeta_comentario(slide, x, y, w, h, encabezado, texto, color_borde=ROJO):
    cont = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    cont.fill.solid(); cont.fill.fore_color.rgb = BLANCO
    cont.line.color.rgb = color_borde; cont.line.width = Pt(1)
    cont.shadow.inherit = False
    hdr_h = Cm(0.6)
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, hdr_h)
    hdr.fill.solid(); hdr.fill.fore_color.rgb = color_borde
    _set_no_line(hdr); hdr.shadow.inherit = False
    tf = hdr.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0.2); tf.margin_top = 0; tf.margin_bottom = 0
    _texto(tf, encabezado, size=10, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    body = slide.shapes.add_textbox(x + Cm(0.15), y + hdr_h + Cm(0.1), w - Cm(0.3), h - hdr_h - Cm(0.2))
    tf2 = body.text_frame; tf2.word_wrap = True
    tf2.margin_left = 0; tf2.margin_top = 0
    _texto(tf2, texto, size=9, color=NEGRO, align=PP_ALIGN.LEFT)
    return cont


def tarjeta_stat(slide, x, y, w, h, valor, etiqueta, color_valor=ROJO, nota=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.08
    box.fill.solid(); box.fill.fore_color.rgb = FONDO_GRIS
    box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD); box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.2); tf.margin_right = Cm(0.2)
    tf.margin_top = Cm(0.12); tf.margin_bottom = Cm(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    p0.text = valor
    for r in p0.runs:
        r.font.size = Pt(19); r.font.bold = True; r.font.color.rgb = color_valor; r.font.name = FUENTE
    p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
    p1.text = etiqueta
    for r in p1.runs:
        r.font.size = Pt(8.5); r.font.color.rgb = GRIS_OSC; r.font.name = FUENTE
    if nota:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        p2.text = nota
        for r in p2.runs:
            r.font.size = Pt(7.5); r.font.italic = True; r.font.color.rgb = GRIS_MED; r.font.name = FUENTE
    return box


def _colorear_serie(serie, color):
    serie.format.fill.solid()
    serie.format.fill.fore_color.rgb = color
    serie.format.line.color.rgb = color


def _colorear_linea(serie, color, dash=None, ancho_pt=2.25):
    serie.format.line.color.rgb = color
    serie.format.line.width = Pt(ancho_pt)
    serie.marker.format.fill.solid()
    serie.marker.format.fill.fore_color.rgb = color
    serie.marker.format.line.color.rgb = color
    if dash:
        ln = serie.format.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': dash}))


def panel_grafico_barras(slide, x, y, w, h, titulo, categorias, series, num_fmt="#,##0"):
    _add_textbox(slide, x, y, w, Cm(0.55), titulo, size=12, bold=True, color=GRIS_OSC)
    chart_data = CategoryChartData()
    chart_data.categories = categorias
    for nombre, valores, _ in series:
        chart_data.add_series(nombre, valores)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y + Cm(0.6), w, h - Cm(0.6), chart_data)
    chart = gf.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    for i, (nombre, valores, color) in enumerate(series):
        _colorear_serie(chart.series[i], color)
    va = chart.value_axis
    va.tick_labels.font.size = Pt(8.5)
    va.format.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    va.tick_labels.number_format = num_fmt
    va.tick_labels.number_format_is_linked = False
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(9)
    ca.format.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    chart.font.size = Pt(9)
    return gf


def panel_grafico_lineas(slide, x, y, w, h, titulo, categorias, series, num_fmt="0%"):
    _add_textbox(slide, x, y, w, Cm(0.55), titulo, size=12, bold=True, color=GRIS_OSC)
    chart_data = CategoryChartData()
    chart_data.categories = categorias
    for s in series:
        chart_data.add_series(s[0], s[1])
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y + Cm(0.6), w, h - Cm(0.6), chart_data)
    chart = gf.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    for i, s in enumerate(series):
        color = s[2]
        dash = s[3] if len(s) > 3 else None
        _colorear_linea(chart.series[i], color, dash=dash)
    va = chart.value_axis
    va.tick_labels.font.size = Pt(8.5)
    va.format.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    va.tick_labels.number_format = num_fmt
    va.tick_labels.number_format_is_linked = False
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(9)
    ca.format.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    chart.font.size = Pt(9)
    return gf


# ---------------------------------------------------------------------------
# CARGA DE DATOS (reutiliza la logica ya validada de generar_kpis.py)
# ---------------------------------------------------------------------------
def cargar_kpis():
    import pandas as pd
    f_det = g.buscar_archivo("detalle", "lote")
    f_cts = g.buscar_archivo("cts", "individual")
    f_ext = g.buscar_archivo("extorno")
    f_abo = g.buscar_archivo("abono", "sueldo")
    f_req = g.buscar_archivo("requerimiento", ext=("csv",))
    if not f_det:
        print("\nERROR: no encuentro 'Detalle de Pago Lote y CTS'.")
        print("Corre primero:  python generar_datos_demo.py")
        sys.exit(1)
    det = g.leer_detalle_pago_lote(f_det)
    cts = g.leer_detalle_cts(f_cts) if f_cts else pd.DataFrame()
    ext = g.leer_extornos(f_ext) if f_ext else pd.DataFrame()
    abo = g.leer_abono(f_abo) if f_abo else pd.DataFrame()
    req = g.leer_requerimientos(f_req) if f_req else pd.DataFrame()
    return g.calcular(det, cts, ext, abo, req)


# ---------------------------------------------------------------------------
# SLIDE · PERFORMANCE KPIs (parte 1: volumen, monto, ritmo de trabajo, errores)
# ---------------------------------------------------------------------------
def construir_slide_kpis_1(slide, K, mes_obj, meses):
    banner_titulo(slide, "Performance KPIs (1/2)",
                  f"Pagos Masivos y CTS · {MESES_LABEL.get(mes_obj, mes_obj)} · Comparativo del periodo")

    resumen = {r["Mes"]: r for r in K["resumen"]}
    cats = [MESES_LABEL.get(m, m) for m in meses]

    y0 = Cm(2.55)
    alto_panel = Cm(6.7)
    ancho_panel = Cm(14.2)
    gap = Cm(0.5)
    x1 = Cm(0.3)
    x2 = x1 + ancho_panel + gap

    # ---- Panel 1: Volumen ----
    k1 = [resumen[m]["K1 Vol Pago"] for m in meses]
    k3 = [resumen[m]["K3 Vol CTS"] for m in meses]
    panel_grafico_barras(slide, x1, y0, ancho_panel, alto_panel - Cm(2.1),
                          "Cuántas transacciones se hicieron cada mes", cats,
                          [("Pago de convenios (planillas)", k1, ROJO), ("Depósitos de CTS", k3, NARANJA)])
    dif = resumen[mes_obj]["K1 Vol Pago"] - resumen[meses[-2]]["K1 Vol Pago"] if len(meses) > 1 else 0
    tarjeta_comentario(slide, x1, y0 + alto_panel - Cm(2.1), ancho_panel, Cm(2.1),
                        "QUÉ SIGNIFICA",
                        f"{MESES_LABEL.get(mes_obj, mes_obj)}: {resumen[mes_obj]['K1 Vol Pago']:,} pagos de convenio "
                        f"({'+' if dif >= 0 else ''}{dif:,} vs. mes anterior). [Dato de ejemplo generado al azar.]")

    # ---- Panel 2: Monto ----
    k2 = [resumen[m]["K2 Monto Pago"] / 1000 for m in meses]
    k4 = [resumen[m]["K4 Monto CTS"] / 1000 for m in meses]
    panel_grafico_barras(slide, x2, y0, ancho_panel, alto_panel - Cm(2.1),
                          "Cuánto dinero se movió cada mes (miles de soles)", cats,
                          [("Pago de convenios", k2, ROJO), ("Depósitos de CTS", k4, NARANJA)])
    tot_mes = resumen[mes_obj]["Total Monto"]
    tarjeta_comentario(slide, x2, y0 + alto_panel - Cm(2.1), ancho_panel, Cm(2.1),
                        "QUÉ SIGNIFICA",
                        f"{MESES_LABEL.get(mes_obj, mes_obj)}: S/ {tot_mes:,.0f} movilizados entre pagos de convenio "
                        f"y CTS, 100% trazable en el sistema.")

    # ---- Panel 3: Ritmo de trabajo del equipo ----
    y1 = y0 + alto_panel + Cm(0.35)
    op_p = [resumen[m]["Ops/dia (persona)"] for m in meses]
    op_c = [resumen[m]["Ops/dia (calend.)"] for m in meses]
    panel_grafico_lineas(slide, x1, y1, ancho_panel, alto_panel - Cm(2.1),
                          "Cuántas transacciones hace cada persona en un día", cats,
                          [("Contando solo días que trabajó cada uno", op_p, ROJO),
                           ("Contando todos los días del mes", op_c, NARANJA)],
                          num_fmt="#,##0")
    tarjeta_comentario(slide, x1, y1 + alto_panel - Cm(2.1), ancho_panel, Cm(2.1),
                        "QUÉ SIGNIFICA",
                        f"{MESES_LABEL.get(mes_obj, mes_obj)}: cada persona atendió ~{resumen[mes_obj]['Ops/dia (persona)']:.0f} "
                        f"transacciones por día trabajado. La brecha con la línea naranja son días sin actividad.")

    # ---- Panel 4: Errores que se revirtieron (extornos) ----
    ext_por_mes = {e["Mes"]: e for e in K["extornos"]}
    meses_ext = [m for m in meses if m in ext_por_mes]
    cats_ext = [MESES_LABEL.get(m, m) for m in meses_ext]
    tasa_imp = [ext_por_mes[m]["Tasa Imputable"] for m in meses_ext]
    tasa_tot = [ext_por_mes[m]["Tasa Total"] for m in meses_ext]
    meta_line = [0.01 for _ in meses_ext]
    panel_grafico_lineas(slide, x2, y1, ancho_panel, alto_panel - Cm(2.1),
                          "Errores que hubo que revertir (meta: menos de 1 de cada 100)", cats_ext,
                          [("Por error del operador", tasa_imp, ROJO),
                           ("Todos los errores (sistema + cliente + operador)", tasa_tot, GRIS_MED),
                           ("Meta: menos de 1%", meta_line, VERDE, "dash")],
                          num_fmt="0.0%")
    ext_mes = ext_por_mes.get(mes_obj)
    if ext_mes:
        txt_ext = (f"{MESES_LABEL.get(mes_obj, mes_obj)}: {ext_mes['Tasa Imputable']*100:.2f} de cada 100 pagos se "
                   f"revirtió por error propio — {ext_mes['Imputables (Error Usuario)']} errores sobre "
                   f"{ext_mes['K1 Pago Lote']:,} pagos hechos.")
    else:
        txt_ext = f"Sin datos de errores (extornos) para {MESES_LABEL.get(mes_obj, mes_obj)}."
    tarjeta_comentario(slide, x2, y1 + alto_panel - Cm(2.1), ancho_panel, Cm(2.1),
                        "QUÉ SIGNIFICA", txt_ext)

    # ---- Franja inferior: tarjetas de indicadores clave ----
    y2 = y1 + alto_panel + Cm(0.3)
    jul = resumen[mes_obj]
    stats = [
        (f"{jul['K5 Ratio Asist.']}x" if jul["K5 Ratio Asist."] != "N/A" else "N/A",
         "Diferencia de carga entre los 2 asistentes", "no debería pasar el doble"),
        (f"{jul['K5a % Analista']*100:.0f}%",
         "Trabajo operativo que aún hace el analista", "debería bajar de 10%"),
        (f"{jul['K7 % Cobert. Pago']*100:.0f}%",
         "Agencias que ya usan Pago en Lote", f"de {g.TOTAL_TIENDAS_RED} agencias en total"),
        (f"{jul['K8 % Cobert. CTS']*100:.0f}%",
         "Agencias que ya usan CTS en Lote", f"de {g.TOTAL_TIENDAS_RED} agencias en total"),
        (f"{jul['K9 Empleadores']}",
         "Empresas distintas atendidas en CTS este mes", None),
    ]
    n = len(stats)
    ancho_total = Cm(28.9)
    gap_s = Cm(0.3)
    ancho_card = (ancho_total - gap_s * (n - 1)) / n
    xs = Cm(0.3)
    for i, (valor, etiqueta, nota) in enumerate(stats):
        tarjeta_stat(slide, xs + i * (ancho_card + gap_s), y2, ancho_card, Cm(2.15),
                     valor, etiqueta, color_valor=ROJO, nota=nota)


# ---------------------------------------------------------------------------
# SLIDE · PERFORMANCE KPIs (parte 2: tiempo estimado)
# ---------------------------------------------------------------------------
def construir_slide_kpis_2(slide, K, mes_obj, meses):
    banner_titulo(slide, "Performance KPIs (2/2)",
                  f"Pagos Masivos y CTS · {MESES_LABEL.get(mes_obj, mes_obj)} · Tiempo estimado de trabajo")

    cats = [MESES_LABEL.get(m, m) for m in meses]

    y0 = Cm(2.55)
    ancho_panel = Cm(28.9)
    alto_panel = Cm(8.0)
    x1 = Cm(0.3)

    tiempo_por_mes = {t["Mes"]: t for t in K["tiempo"]}
    hp = [tiempo_por_mes[m]["Horas Pago Lote"] for m in meses]
    hc = [tiempo_por_mes[m]["Horas CTS"] for m in meses]
    panel_grafico_barras(slide, x1, y0, ancho_panel, alto_panel - Cm(2.1),
                          "Horas de trabajo que representan estas transacciones (estimado)", cats,
                          [("Pago de convenios", hp, ROJO), ("Depósitos de CTS", hc, NARANJA)],
                          num_fmt="#,##0.0")
    t_mes = tiempo_por_mes.get(mes_obj)
    tarjeta_comentario(slide, x1, y0 + alto_panel - Cm(2.1), ancho_panel, Cm(2.1),
                        "QUÉ SIGNIFICA",
                        f"Es un estimado, no un cronómetro real: se calcula con un tiempo promedio por transacción "
                        f"({g.MIN_POR_PAGO_LOTE} min por pago de convenio, {g.MIN_POR_CTS} min por CTS). En "
                        f"{MESES_LABEL.get(mes_obj, mes_obj)} equivale a {t_mes['Horas Totales Estim.']:.0f} horas "
                        f"de trabajo, unos {t_mes['Equiv. Dias (8h)']:.1f} días completos de una persona sola.")

    _add_textbox(slide, Cm(0.3), y0 + alto_panel + Cm(0.15), Cm(28.9), Cm(1.2),
                 "Nota: los tiempos son un estimado con supuestos fijos (no cronometrado).",
                 size=9, color=GRIS_MED)


# ---------------------------------------------------------------------------
# SLIDE · GESTION OPERATIVA (datos de ejemplo, no reales)
# ---------------------------------------------------------------------------
def construir_slide_gestion(slide, K, mes_obj):
    banner_titulo(slide, "Gestión Operativa", f"Pagos Masivos y CTS · {MESES_LABEL.get(mes_obj, mes_obj)}")

    y = Cm(2.55)

    # ---- Logro del mes (EJEMPLO) ----
    logro = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(0.3), y, Cm(28.9), Cm(1.5))
    logro.adjustments[0] = 0.15
    logro.fill.solid(); logro.fill.fore_color.rgb = ROJO
    _set_no_line(logro); logro.shadow.inherit = False
    tf = logro.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0.5); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    p.text = "🏆  LOGRO DEL MES (EJEMPLO):  Salida a producción de una nueva mejora del proceso"
    for r in p.runs:
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = BLANCO; r.font.name = FUENTE

    y += Cm(1.8)

    # ---- Mejoras solicitadas (EJEMPLO — reemplaza por tus propias solicitudes) ----
    _add_textbox(slide, Cm(0.3), y, Cm(14), Cm(0.6),
                 "Mejoras solicitadas este mes (datos de ejemplo)", size=13, bold=True, color=GRIS_OSC)
    y_tabla = y + Cm(0.65)
    filas_hc = [
        ("REQ-001", "01-mes", "Mejorar un reporte operativo",
         "Ejemplo: agregar campos para poder reportar por ubicación.", "En coordinación"),
        ("REQ-002", "02-mes", "Agilizar un proceso interno",
         "Ejemplo: quitar pasos redundantes y agregar validaciones.", "2 semanas est."),
        ("REQ-003", "02-mes", "Automatizar una tarea manual",
         "Ejemplo: que el sistema organice archivos automáticamente.", "4 semanas est."),
        ("REQ-004", "09-mes", "Reforzar un control de riesgo",
         "Ejemplo: que las alertas lleguen también al equipo responsable.", "En coordinación"),
    ]
    col_w = [Cm(2.3), Cm(1.6), Cm(6.0), Cm(15.5), Cm(3.1)]
    headers = ["N° Ref.", "Fecha", "Qué se pidió", "En qué consiste", "Estado"]
    tabla_h = Cm(0.6) + Cm(1.05) * len(filas_hc)
    tabla_shape = slide.shapes.add_table(len(filas_hc) + 1, len(headers), Cm(0.3), y_tabla,
                                          sum(col_w, Cm(0)), tabla_h)
    tabla = tabla_shape.table
    for i, w in enumerate(col_w):
        tabla.columns[i].width = w
    tabla.rows[0].height = Cm(0.6)
    for j, htxt in enumerate(headers):
        cell = tabla.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = ROJO
        cell.text = htxt
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(10); r.font.color.rgb = BLANCO; r.font.name = FUENTE
    for i, fila in enumerate(filas_hc, start=1):
        tabla.rows[i].height = Cm(1.05)
        for j, val in enumerate(fila):
            cell = tabla.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = FONDO_GRIS if i % 2 == 0 else BLANCO
            cell.text = val
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Cm(0.15); cell.margin_right = Cm(0.15)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if j in (0, 1, 4) else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(9); r.font.color.rgb = NEGRO; r.font.name = FUENTE
                    r.font.bold = (j == 0)

    y_buzon = y_tabla + tabla_h + Cm(0.4)

    # ---- Buzon de Requerimientos ----
    _add_textbox(slide, Cm(0.3), y_buzon, Cm(14), Cm(0.6), "Buzón de Requerimientos",
                 size=13, bold=True, color=GRIS_OSC)
    y_buzon_body = y_buzon + Cm(0.65)

    buzon_tipo = K.get("buzon_tipo") or []
    buzon_asig = K.get("buzon_asig") or []
    total_sol = K.get("buzon_total") or sum(b["Cantidad"] for b in buzon_tipo)
    top3 = buzon_tipo[:3]

    tarjeta_stat(slide, Cm(0.3), y_buzon_body, Cm(4.6), Cm(2.6),
                 f"{total_sol:,}", "Solicitudes recibidas (acumulado)", color_valor=ROJO)

    x_top = Cm(5.2)
    _add_textbox(slide, x_top, y_buzon_body, Cm(9.5), Cm(0.5), "Por qué escriben más los clientes", size=10.5, bold=True, color=GRIS_OSC)
    yy = y_buzon_body + Cm(0.55)
    for item in top3:
        barra_w = Cm(5.5) * (item["Cantidad"] / top3[0]["Cantidad"]) if top3 else Cm(0)
        barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_top, yy, barra_w, Cm(0.45))
        barra.fill.solid(); barra.fill.fore_color.rgb = NARANJA
        _set_no_line(barra); barra.shadow.inherit = False
        _add_textbox(slide, x_top + barra_w + Cm(0.15), yy - Cm(0.02), Cm(6.5), Cm(0.5),
                     f"{item['Tipo de Solicitud'].title()} — {item['Cantidad']} ({item['% del Total']*100:.0f}%)",
                     size=9.5, color=NEGRO)
        yy += Cm(0.6)

    x_asig = Cm(15.5)
    _add_textbox(slide, x_asig, y_buzon_body, Cm(9), Cm(0.5), "Quién atendió cada solicitud", size=10.5, bold=True, color=GRIS_OSC)
    yy2 = y_buzon_body + Cm(0.55)
    top_asig = buzon_asig[:4]
    for item in top_asig:
        barra_w = Cm(5.0) * (item["Cantidad"] / top_asig[0]["Cantidad"]) if top_asig else Cm(0)
        barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_asig, yy2, barra_w, Cm(0.4))
        barra.fill.solid(); barra.fill.fore_color.rgb = ROJO
        _set_no_line(barra); barra.shadow.inherit = False
        _add_textbox(slide, x_asig + barra_w + Cm(0.15), yy2 - Cm(0.02), Cm(5.5), Cm(0.45),
                     f"{item['Asignado a'].title()} — {item['Cantidad']}",
                     size=9, color=NEGRO)
        yy2 += Cm(0.52)

    y_alerta = max(yy, yy2) + Cm(0.25)
    alerta = K.get("buzon_alerta")
    if alerta:
        caja_al = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(0.3), y_alerta, Cm(28.9), Cm(1.4))
        caja_al.adjustments[0] = 0.15
        caja_al.fill.solid(); caja_al.fill.fore_color.rgb = RGBColor(0xFF, 0xEA, 0xEA)
        caja_al.line.color.rgb = ROJO; caja_al.line.width = Pt(1)
        caja_al.shadow.inherit = False
        tf = caja_al.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Cm(0.5); tf.margin_right = Cm(0.5)
        sin_fecha_txt = alerta.replace("BRECHA: ", "").split(" -> ")[0]
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        p.text = f"⚠ Punto de atención: {sin_fecha_txt}."
        for r in p.runs:
            r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = ROJO; r.font.name = FUENTE
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        p2.text = "Sin esa fecha no se puede medir si el buzón responde rápido o lento."
        for r in p2.runs:
            r.font.size = Pt(10); r.font.color.rgb = NEGRO; r.font.name = FUENTE


# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------
def main():
    print("Cargando KPIs (datos disponibles en esta carpeta)...")
    K = cargar_kpis()
    meses = K["meses"]
    if not meses:
        print("ERROR: no se detectaron meses con datos."); sys.exit(1)
    mes_obj = meses[-1]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    layout_blanco = prs.slide_layouts[6]  # layout en blanco

    # Limpia metadatos residuales de la plantilla interna de python-pptx
    # (autor, ultima modificacion, etc.) para no dejar rastros ajenos.
    cp = prs.core_properties
    cp.author = ""
    cp.last_modified_by = ""
    cp.title = "Informe de KPIs — Pagos Masivos y CTS (demo)"
    cp.subject = ""
    cp.comments = ""

    print("Construyendo Slide 1 · Performance KPIs (1/2)...")
    s1 = prs.slides.add_slide(layout_blanco)
    construir_slide_kpis_1(s1, K, mes_obj, meses)

    print("Construyendo Slide 2 · Performance KPIs (2/2)...")
    s2 = prs.slides.add_slide(layout_blanco)
    construir_slide_kpis_2(s2, K, mes_obj, meses)

    print("Construyendo Slide 3 · Gestion Operativa (datos de ejemplo)...")
    s3 = prs.slides.add_slide(layout_blanco)
    construir_slide_gestion(s3, K, mes_obj)

    prs.save(SALIDA)
    print("Guardado:", os.path.basename(SALIDA), f"({len(prs.slides)} diapositivas)")


if __name__ == "__main__":
    main()
